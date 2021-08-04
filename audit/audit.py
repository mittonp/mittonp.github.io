import csv
from pptx import Presentation
from datetime import date
import time

with open('C:/Users/mitto/2021.07.28.04.56.14/response_codes_redirection_(3xx).csv') as csv_file:
	csv_reader = csv.reader(csv_file, delimiter=',')
	line_count = 0
	for row in csv_reader:
		line_count += 1
print(f'There are {line_count-1} 3xx response codes.')

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Vine Digital Automated Audit"
subtitle.text = "Generated on " + date.today().strftime("%#d/%m/%y")
# slide.placeholders[2].text = "digitalitinerant.com"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "3xx responses"
slide.placeholders[1].text = str(line_count -1)

with open('C:/Users/mitto/2021.07.28.04.56.14/client_error_(4xx)_inlinks.csv') as csv_file:
	csv_reader = csv.reader(csv_file, delimiter=',')
	line_count = 0
	for row in csv_reader:
		line_count += 1
		print (row)
print(f'There are {line_count-1} 4xx response codes.')

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "4xx responses"
slide.placeholders[1].text = str(line_count -1)

prs.save('test.pptx')